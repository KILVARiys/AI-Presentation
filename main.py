import os
from transformers import pipeline
from pptx import Presentation
from pptx.util import Inches, Pt

# Указываем модель и ревизию
model_name = "facebook/bart-large-cnn"
revision = "main"

# Функция для генерации сводки
def generate_summary(input_text):
    summarizer = pipeline("summarization", model=model_name, revision=revision)
    summary = summarizer(input_text, max_length=80, min_length=20, do_sample=False)
    return summary[0]['summary_text']

# Функция для создания презентации
def create_presentation(summary_text, image_paths):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Пустой слайд

    # Добавление заголовка
    title = slide.shapes.title
    title.text = "Заголовок"

    # Добавление текста
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)  # Увеличим высоту, чтобы уместить текст
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame

    # Разделяем текст на предложения и добавляем их в новый параграф
    for sentence in summary_text.split('.'):
        sentence = sentence.strip()  # Убираем лишние пробелы
        if sentence:  # Проверяем, что предложение не пустое
            p = text_frame.add_paragraph()
            p.text = sentence + '.'  # Добавляем точку в конец предложения
            p.font.size = Pt(18)  # Устанавливаем размер шрифта

    # Добавление изображений
    for idx, image_path in enumerate(image_paths):
        left = Inches(1)  # Положение по горизонтали
        top = Inches(3 + idx * 2)  # Положение по вертикали
        slide.shapes.add_picture(image_path, left, top, width=Inches(5))

    # Путь к папке Загрузки
    downloads_path = os.path.join(os.path.expanduser("~"), "Downloads", "presentation.pptx")
    prs.save(downloads_path)

    return downloads_path  # Возвращаем путь к сохраненному файлу

# Основной код
if __name__ == "__main__":
    # Входной текст
    input_text = "Basilio the cat checks the presentation for functionality."

    # Генерация сводки
    summary = generate_summary(input_text)
    print("Сводка:", summary)

    # Пути к изображениям (замени на реальные пути к изображениям)
    image_paths = [
        "D:/image1.png",
    ]

    # Создание презентации и получение пути к файлу
    saved_path = create_presentation(summary, image_paths)
    print(f"Презентация сохранена в: {saved_path}")
